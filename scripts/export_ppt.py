import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_ppt_from_yaml(data_file, output_file):
    try:
        with open(data_file, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
    except Exception as e:
        print(f"❌ Error loading YAML {data_file}: {e}")
        return

    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = data.get('title', 'Untitled')
    if 'subtitle' in data:
        slide.placeholders[1].text = data['subtitle']

    # Objectives Slide
    if 'objectives' in data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "학습 목표 🎯"
        tf = slide.placeholders[1].text_frame
        for obj in data['objectives']:
            p = tf.add_paragraph()
            p.text = obj
            p.level = 0

    # Sections
    for section in data.get('sections', []):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = section.get('title', '')
        
        tf = slide.placeholders[1].text_frame
        
        if 'content' in section:
            # Simple HTML tag stripping
            clean_text = section['content'].replace('<br>', '\n').replace('<b>', '').replace('</b>', '').strip()
            # Remove div tags roughly
            clean_text = clean_text.split('<div')[0] 
            
            p = tf.add_paragraph()
            p.text = clean_text[:500]

        if 'columns' in section:
            for col in section['columns']:
                if 'title' in col:
                    p = tf.add_paragraph()
                    p.text = f"▶ {col['title']}"
                    p.font.bold = True
                if 'list' in col:
                    for item in col['list']:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 1
                if 'text' in col:
                    p = tf.add_paragraph()
                    p.text = col['text'].replace('<b>', '').replace('</b>', '')
                    p.level = 1

    # Save
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    prs.save(output_file)
    print(f"✅ Generated PPT: {os.path.basename(output_file)}")

if __name__ == "__main__":
    pass
