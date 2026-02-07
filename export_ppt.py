import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_ppt_from_yaml(data_file, output_file):
    # 1. Load Data
    with open(data_file, 'r', encoding='utf-8') as f:
        data = yaml.safe_load(f)

    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = data['title']
    slide.placeholders[1].text = data['subtitle']

    # Objectives Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "학습 목표 🎯"
    tf = slide.placeholders[1].text_frame
    for obj in data['objectives']:
        p = tf.add_paragraph()
        p.text = obj
        p.level = 0

    # Sections Slides
    for section in data.get('sections', []):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = section['title']
        
        tf = slide.placeholders[1].text_frame
        
        # Add content if exists
        if 'content' in section:
            # Strip simple HTML-like tags for PPT
            clean_text = section['content'].replace('<br>', '\n').replace('<b>', '').replace('</b>', '').strip()
            # Simple check to avoid too long text
            p = tf.add_paragraph()
            p.text = clean_text[:500] 

        # Add column data as bullet points
        if 'columns' in section:
            for col in section['columns']:
                if 'title' in col:
                    p = tf.add_paragraph()
                    p.text = f"▶ {col['title']}"
                    p.font.bold = True
                if 'list' in col:
                    for item in col['list']:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 1
                if 'text' in col:
                    p = tf.add_paragraph()
                    p.text = col['text'].replace('<b>', '').replace('</b>', '')
                    p.level = 1

    # Summary Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "핵심 정리 📝"
    tf = slide.placeholders[1].text_frame
    for item in data['summary']:
        p = tf.add_paragraph()
        p.text = item.replace('<b>', '').replace('</b>', '')
        p.level = 0

    # Save
    prs.save(output_file)
    print(f"✅ PowerPoint generated: {output_file}")

if __name__ == "__main__":
    os.makedirs('Project-SCOPE/output', exist_ok=True)
    create_ppt_from_yaml(
        'Project-SCOPE/data/02_fertilization.yaml',
        'Project-SCOPE/output/02_fertilization_lesson.pptx'
    )
